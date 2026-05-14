# -*- coding: utf-8 -*-
"""
PPT 报告生成器（企业级美化版）
==============================

使用 python-pptx 自动生成现代风格的经营分析演示文稿。
设计原则：简洁、专业、信息层级清晰。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# 企业级配色方案（取自现代 BI 设计系统）
# ---------------------------------------------------------------------------
# 主色调：深蓝 -> 渐变过渡
_C_PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)       # 深海蓝（标题栏、背景）
_C_SECONDARY = RGBColor(0x2C, 0x3E, 0x6B)      # 中蓝（辅助背景）
_C_ACCENT = RGBColor(0x3B, 0x82, 0xF6)         # 明亮蓝（强调色）
_C_ACCENT_LIGHT = RGBColor(0x93, 0xC5, 0xFD)   # 淡蓝
_C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_C_OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)      # 页面背景
_C_LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)     # 分割线 / 装饰
_C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)       # 正文
_C_GRAY_TEXT = RGBColor(0x64, 0x74, 0x8B)       # 辅助文字

# 风险语义色
_C_RED = RGBColor(0xEF, 0x44, 0x44)
_C_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
_C_GREEN = RGBColor(0x10, 0xB9, 0x81)

# 数据色板（图表配色）
_DATA_COLORS = [
    RGBColor(0x3B, 0x82, 0xF6),
    RGBColor(0x10, 0xB9, 0x81),
    RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0xEF, 0x44, 0x44),
    RGBColor(0x8B, 0x5C, 0xF6),
    RGBColor(0xEC, 0x48, 0x99),
]

# ---------------------------------------------------------------------------
# 尺寸常量
# ---------------------------------------------------------------------------
_SW = Inches(13.33)       # slide width (16:9)
_SH = Inches(7.5)         # slide height
_MARGIN_L = Inches(0.8)
_MARGIN_R = Inches(0.8)
_CONTENT_W = _SW - _MARGIN_L - _MARGIN_R
_TITLE_BAR_H = Inches(1.1)


class PPTGenerator:
    """
    企业级 PPT 报告生成器（美化版）。

    用法::

        gen = PPTGenerator()
        gen.generate(df, categories, insights, chart_paths)
    """

    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = _SW
        self.prs.slide_height = _SH

    # ==================================================================
    # 主入口
    # ==================================================================

    def generate(
        self,
        df: Any,
        categories: dict[str, Any],
        insights: str,
        chart_paths: dict[str, str],
        output_path: str | Path = "reports/enterprise_report.pptx",
    ) -> str:
        """生成完整 PPT。"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 页码计数器（从 0 开始，_page_number 先 ++ 再显示）
        self._page_num = 0
        self._df = df
        self._categories = categories
        self._chart_paths = chart_paths

        total = len(df)
        high_risk = categories.get("high_risk", [])
        high_value = categories.get("high_value", [])
        growth = categories.get("growth", [])

        self._add_cover(total, high_risk, high_value, growth)
        self._add_toc(total, high_risk, high_value, growth)
        self._add_executive_summary(total, high_risk, high_value, growth, df, insights)
        self._add_risk_overview(chart_paths.get("risk_pie", ""), df)
        self._add_high_risk_detail(high_risk)
        self._add_high_value_slide(high_value)
        self._add_charts_grid(chart_paths)
        self._add_ai_overview()
        self._add_ai_insights(insights)
        self._add_action_plan(high_risk, insights)
        self._add_closing()

        self.prs.save(str(output_path))
        logger.info("PPT 已保存: %s", output_path)
        return str(output_path)

    # ==================================================================
    # 幻灯片基本构建块
    # ==================================================================

    def _new_slide(self) -> Any:
        """添加空白幻灯片。"""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _page_bg(self, slide: Any, color: RGBColor = _C_OFF_WHITE) -> None:
        """设置页面纯色背景。"""
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = color

    def _title_bar(self, slide: Any, title: str, subtitle: str = "") -> None:
        """统一的顶部标题栏（深蓝底色 + 左侧 accent 条）。"""
        # 深蓝横条
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SW, _TITLE_BAR_H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _C_PRIMARY
        bar.line.fill.background()

        # 左侧 accent 竖条
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), _TITLE_BAR_H, Inches(0.08), Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = _C_ACCENT
        accent.line.fill.background()

        # 标题文字
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = _C_WHITE
        p.font.bold = True

        if subtitle:
            tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11), Inches(0.35))
            tf2 = tb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(12)
            p2.font.color.rgb = _C_ACCENT_LIGHT

    def _footer_line(self, slide: Any) -> None:
        """底部装饰线。"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), _CONTENT_W, Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _C_LIGHT_GRAY
        line.line.fill.background()

    def _page_number(self, slide: Any) -> None:
        """右下角页码（自动递增）。"""
        self._page_num += 1
        tb = slide.shapes.add_textbox(Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = str(self._page_num)
        p.font.size = Pt(9)
        p.font.color.rgb = _C_GRAY_TEXT
        p.alignment = PP_ALIGN.RIGHT

    def _body_box(
        self, slide: Any, left: float, top: float,
        width: float, height: float, text: str,
        font_size: int = 14, color: RGBColor = _C_DARK_TEXT,
        bold: bool = False, line_spacing: float = 1.5,
    ) -> None:
        """带 word_wrap 的正文文本框。"""
        tb = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.strip().split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.space_after = Pt(4)

    def _card(
        self, slide: Any, left: float, top: float,
        width: float, height: float,
        title: str, value: str, accent_color: RGBColor = _C_ACCENT,
    ) -> None:
        """卡片式 KPI 组件。"""
        # 卡片背景（白色圆角矩形 — PPT 用矩形模拟）
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _C_WHITE
        card.line.color.rgb = _C_LIGHT_GRAY
        card.line.width = Pt(0.5)
        # 阴影效果（浅灰偏移）
        shadow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left + 0.03), Inches(top + 0.03), Inches(width), Inches(height),
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shadow.line.fill.background()
        # 把 card 移到最前（z-order）—— python-pptx 按添加顺序，所以先加 shadow 再加 card

        # 顶部 accent 条
        accent_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.06),
        )
        accent_strip.fill.solid()
        accent_strip.fill.fore_color.rgb = accent_color
        accent_strip.line.fill.background()

        # 数值
        tb_v = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.35), Inches(width - 0.4), Inches(0.5),
        )
        tf_v = tb_v.text_frame
        p_v = tf_v.paragraphs[0]
        p_v.text = value
        p_v.font.size = Pt(32)
        p_v.font.color.rgb = _C_DARK_TEXT
        p_v.font.bold = True

        # 标签
        tb_l = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.85), Inches(width - 0.4), Inches(0.3),
        )
        tf_l = tb_l.text_frame
        p_l = tf_l.paragraphs[0]
        p_l.text = title
        p_l.font.size = Pt(12)
        p_l.font.color.rgb = _C_GRAY_TEXT

    def _risk_badge(self, text: str, level: str) -> RGBColor:
        """根据风险等级返回颜色。"""
        if level == "高":
            return _C_RED
        elif level == "中":
            return _C_ORANGE
        return _C_GREEN

    # ==================================================================
    # AI 洞察解析与渲染辅助方法
    # ==================================================================

    def _detect_section_type(self, title: str) -> str:
        """根据标题检测 section 类型。"""
        if "经营概况" in title or "概况" in title:
            return "overview"
        if "风险" in title:
            return "risk"
        if "策略" in title or "分层" in title or "经营策略" in title:
            return "strategy"
        if "客户清单" in title or "重点客户" in title:
            return "list"
        if "行动" in title or "计划" in title:
            return "action"
        return "default"

    def _parse_insight_sections(self, insights: str) -> list[tuple[str, str]]:
        """将 Markdown 按 ## 标题解析为 (标题, 内容) 章节列表。"""
        sections: list[tuple[str, str]] = []
        current_title = "综合洞察"
        current_lines: list[str] = []
        for line in insights.split("\n"):
            if line.startswith("## "):
                if current_lines:
                    sections.append((current_title, "\n".join(current_lines).strip()))
                current_title = line[3:].strip()
                current_lines = []
            else:
                current_lines.append(line)
        if current_lines:
            sections.append((current_title, "\n".join(current_lines).strip()))
        sections = [s for s in sections if len(s[1]) > 10]
        return sections

    def _merge_wrapper_sections(self, sections: list[tuple[str, str]]) -> list[tuple[str, str]]:
        """将"综合洞察/AI 经营洞察"包裹段的内容合并到下一节，避免重复幻灯片。"""
        if len(sections) < 2:
            return sections
        first_title = sections[0][0]
        if any(kw in first_title for kw in ("洞察", "综合")):
            wrapper_content = sections[0][1]
            main_title, main_content = sections[1]
            sections[1] = (main_title, wrapper_content + "\n\n" + main_content)
            sections.pop(0)
        return sections

    def _clean_markdown_text(self, text: str) -> str:
        """去除 Markdown 格式标记，保留内容。"""
        import re
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        lines = []
        for line in text.split("\n"):
            sl = line.strip()
            if sl.startswith("- ") or sl.startswith("* "):
                sl = "• " + sl[2:]
            lines.append(sl)
        return "\n".join(lines)

    def _extract_strategies(self, content: str) -> dict:
        """从 AI 内容中提取各策略描述。"""
        strategies: dict[str, str] = {}
        lines = content.split("\n")
        current_key: str | None = None
        for line in lines:
            sl = line.strip()
            if "高风险" in sl and any(kw in sl for kw in ("策略", "挽回", "客户")):
                current_key = "high_risk"
            elif "高价值" in sl and any(kw in sl for kw in ("策略", "维护", "客户")):
                current_key = "high_value"
            elif "增长" in sl and any(kw in sl for kw in ("策略", "潜力", "客户")):
                current_key = "growth"
            elif "行业" in sl and any(kw in sl for kw in ("策略", "专项")):
                current_key = "industry"
            elif sl.startswith("- ") and current_key and current_key not in strategies:
                strategies[current_key] = sl[2:]
        defaults = {
            "high_risk": "高层拜访 + 紧急续约谈判，制定个性化挽回方案",
            "high_value": "专属 CSM 配备 + 季度 QBR 复盘 + 提前 90 天续约提醒",
            "growth": "CSM 专项辅导 + 产品深度培训 + 行业案例分享",
            "industry": "行业痛点调研 + 功能定制 + 专项服务小组",
        }
        for k, v in defaults.items():
            strategies.setdefault(k, v)
        return strategies

    def _extract_action_items(self, content: str) -> list[str]:
        """从 AI 内容中提取行动项。"""
        items: list[str] = []
        for line in content.split("\n"):
            sl = line.strip()
            if not sl or "##" in sl:
                continue
            if sl[0].isdigit() and ". " in sl[:4]:
                parts = sl.split(". ", 1)
                if len(parts) > 1:
                    items.append(parts[1].replace("**", "").replace("*", ""))
            elif sl.startswith("- ") or sl.startswith("* "):
                items.append(sl[2:].replace("**", "").replace("*", ""))
        return items[:5]

    def _parse_customers_from_content(self, content: str) -> list[dict]:
        """从 AI 文本中尝试解析客户列表（备用 fallback）。"""
        import re
        customers: list[dict] = []
        current: dict = {}
        for line in content.split("\n"):
            sl = line.strip()
            m = re.match(r'^\d+[.、]\s*\*{0,2}(.+?)\*{0,2}\s*[（(](.+?)[）)]', sl)
            if m:
                if current:
                    customers.append(current)
                current = {
                    "customer_name": m.group(1).strip(),
                    "industry": m.group(2).strip(),
                }
            elif sl.startswith("- ") and current:
                key_val = sl[2:]
                if "风险" in key_val:
                    current["risk_level"] = key_val.split("：")[-1] if "：" in key_val else "高"
        if current:
            customers.append(current)
        return customers

    # ==================================================================
    # 各页幻灯片
    # ==================================================================

    # ---- 1. 封面 ----

    def _add_cover(self, total: int, high_risk: list, high_value: list, growth: list) -> None:
        """现代风格封面。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_PRIMARY)

        # 左上装饰网格线（几何装饰）
        for i in range(6):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5 + i * 0.3), Inches(0.5), Inches(0.01), Inches(1.5),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = _C_SECONDARY
            line.line.fill.background()

        # 主标题区域 - 左侧竖线 accent
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(2.2), Inches(0.08), Inches(1.8),
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = _C_ACCENT
        accent_line.line.fill.background()

        # 主标题
        tb = slide.shapes.add_textbox(Inches(1.6), Inches(2.3), Inches(10), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "企业经营分析报告"
        p.font.size = Pt(46)
        p.font.color.rgb = _C_WHITE
        p.font.bold = True

        # 副标题
        tb2 = slide.shapes.add_textbox(Inches(1.6), Inches(3.5), Inches(10), Inches(0.6))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "Enterprise Business Intelligence Report"
        p2.font.size = Pt(18)
        p2.font.color.rgb = _C_ACCENT_LIGHT

        # 分割线
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.6), Inches(4.3), Inches(3), Inches(0.02),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = _C_ACCENT
        div.line.fill.background()

        # 日期 + 来源
        import pandas as pd
        tb3 = slide.shapes.add_textbox(Inches(1.6), Inches(4.5), Inches(10), Inches(0.4))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"生成日期：{pd.Timestamp.now().strftime('%Y 年 %m 月 %d 日')}  |  AI-Customer-Success-Copilot"
        p3.font.size = Pt(12)
        p3.font.color.rgb = _C_GRAY_TEXT

        # 底部四个 KPI 卡片
        kpi_data = [
            ("总客户数", str(total), _C_ACCENT),
            ("高风险", str(len(high_risk)), _C_RED),
            ("高价值", str(len(high_value)), _C_GREEN),
            ("增长潜力", str(len(growth)), _C_ORANGE),
        ]
        for i, (label, value, color) in enumerate(kpi_data):
            x = 1.2 + i * 2.7
            self._card(slide, x, 5.5, 2.3, 1.3, label, value, color)

    # ---- 2. 目录 ----

    def _add_toc(self, total: int, high_risk: list, high_value: list, growth: list) -> None:
        """目录页。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)

        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(5), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "目录 CONTENTS"
        p.font.size = Pt(24)
        p.font.color.rgb = _C_PRIMARY
        p.font.bold = True

        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.03),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = _C_ACCENT
        div.line.fill.background()

        items = [
            ("01", "执行摘要", "整体经营状况一览"),
            ("02", "风险总览", "客户风险分布与统计"),
            ("03", "高风险客户详情", "需紧急干预的客户列表"),
            ("04", "高价值客户", "核心客户经营状况"),
            ("05", "行业与区域分析", "多维度数据图表"),
            ("06", "AI 经营洞察", "DeepSeek 智能分析"),
            ("07", "行动计划", "分级跟进策略"),
        ]
        for i, (num, title, desc) in enumerate(items):
            y = 1.6 + i * 0.75
            # 编号圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = _C_PRIMARY
            circle.line.fill.background()
            tf_c = circle.text_frame
            tf_c.paragraphs[0].text = num
            tf_c.paragraphs[0].font.size = Pt(12)
            tf_c.paragraphs[0].font.color.rgb = _C_WHITE
            tf_c.paragraphs[0].font.bold = True
            tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf_c.word_wrap = False

            # 标题
            tb_t = slide.shapes.add_textbox(Inches(1.5), Inches(y - 0.05), Inches(5), Inches(0.35))
            tf_t = tb_t.text_frame
            tf_t.paragraphs[0].text = title
            tf_t.paragraphs[0].font.size = Pt(16)
            tf_t.paragraphs[0].font.color.rgb = _C_DARK_TEXT
            tf_t.paragraphs[0].font.bold = True

            # 描述
            tb_d = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.25), Inches(5), Inches(0.3))
            tf_d = tb_d.text_frame
            tf_d.paragraphs[0].text = desc
            tf_d.paragraphs[0].font.size = Pt(11)
            tf_d.paragraphs[0].font.color.rgb = _C_GRAY_TEXT

            # 分割虚线
            if i < len(items) - 1:
                dot_line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.7), Inches(6), Inches(0.005),
                )
                dot_line.fill.solid()
                dot_line.fill.fore_color.rgb = _C_LIGHT_GRAY
                dot_line.line.fill.background()

    # ---- 3. 执行摘要 ----

    def _add_executive_summary(
        self, total: int, high_risk: list, high_value: list, growth: list, df: Any,
        ai_insights: str = "",
    ) -> None:
        """执行摘要页。

        当 ai_insights 包含 AI 生成的有效内容时，右侧面板展示 AI 摘要。
        """
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(slide, "执行摘要", "Executive Summary")

        import pandas as pd
        avg_gmv = df["monthly_gmv"].mean()
        risk_counts = df["risk_level"].value_counts().to_dict()
        industries = df["industry"].nunique()
        regions = df["region"].nunique()

        # 顶部 KPI 卡片行
        kpis = [
            ("总客户数", str(total), _C_ACCENT),
            ("覆盖行业", f"{industries} 个", _C_GREEN),
            ("覆盖区域", f"{regions} 个", _C_ORANGE),
            ("月均 GMV", f"¥{avg_gmv:,.0f}", _C_PRIMARY),
        ]
        for i, (label, value, color) in enumerate(kpis):
            x = 0.8 + i * 3.1
            self._card(slide, x, 1.4, 2.8, 1.2, label, value, color)

        # 风险分布 + 注释
        risk_text = (
            "▎客户风险分布\n\n"
            f"🔴  高风险  {risk_counts.get('高', 0)} 个  ——  需立即启动挽回流程\n"
            f"🟡  中风险  {risk_counts.get('中', 0)} 个  ——  需关注并制定改善计划\n"
            f"🟢  低风险  {risk_counts.get('低', 0)} 个  ——  正常维护\n\n"
            f"高价值客户：{len(high_value)} 个  |  增长潜力客户：{len(growth)} 个"
        )
        self._body_box(slide, 0.8, 3.0, 5.5, 3.5, risk_text, font_size=14)

        # 右侧：AI 洞察内容（有 AI 时展示 AI 摘要，否则展示模板）
        if ai_insights and not ai_insights.startswith("（AI 洞察") and not ai_insights.startswith("## AI 经营洞察"):
            # 提取 AI 摘要的前 300 字作为核心发现
            short = ai_insights[:300].strip()
            # 移除 Markdown 标题标记，简化显示
            short = short.replace("## ", "").replace("**", "")
            insight_text = f"▎AI 核心发现\n\n{short}\n\n……（完整内容见 AI 洞察页）"
        else:
            insight_text = (
                "▎核心发现\n\n"
                f"• 高风险客户占比 {risk_counts.get('高', 0) / max(total, 1) * 100:.0f}%，"
                f"需优先处理\n"
                f"• 高价值客户贡献稳定，重点维系\n"
                f"• 增长潜力客户有 uplift 空间\n\n"
                "▎建议\n\n"
                "• 本周内完成高风险客户 1 对 1 沟通\n"
                "• 高价值客户安排 QBR 季度复盘\n"
                "• 增长客户 CSM 制定专项护航计划"
            )
        self._body_box(slide, 7.0, 3.0, 5.5, 3.5, insight_text, font_size=14)

        self._footer_line(slide)
        self._page_number(slide)

    # ---- 4. 风险总览 ----

    def _add_risk_overview(self, chart_path: str, df: Any) -> None:
        """客户风险总览页（图表 + 定义）。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(slide, "客户风险总览", "Risk Overview")

        # 左侧：饼图
        if chart_path:
            slide.shapes.add_picture(
                chart_path, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5),
            )

        # 右侧：风险定义 + 统计
        risk_counts = df["risk_level"].value_counts().to_dict() if hasattr(df, 'value_counts') else {}

        stats = [
            ("高风险", risk_counts.get('高', 0), _C_RED,
             "续费逾期 / 登录 < 5d / 投诉 > 8 次\n48 小时内必须触发挽回流程"),
            ("中风险", risk_counts.get('中', 0), _C_ORANGE,
             "部分指标不达标\n需 CSM 本周内制定改善计划"),
            ("低风险", risk_counts.get('低', 0), _C_GREEN,
             "各项指标正常\n按季度常规回访即可"),
        ]
        for i, (label, count, color, desc) in enumerate(stats):
            y = 1.6 + i * 1.8
            # 色块
            dot = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(y), Inches(0.06), Inches(0.5),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            # 等级标题
            tb = slide.shapes.add_textbox(Inches(7.1), Inches(y - 0.05), Inches(2.5), Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label}  {count} 个"
            p.font.size = Pt(16)
            p.font.color.rgb = _C_DARK_TEXT
            p.font.bold = True

            # 描述
            tb2 = slide.shapes.add_textbox(Inches(7.1), Inches(y + 0.35), Inches(5.5), Inches(1.2))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = desc
            p2.font.size = Pt(11)
            p2.font.color.rgb = _C_GRAY_TEXT

        self._footer_line(slide)
        self._page_number(slide)

    # ---- 5. 高风险客户详情 ----

    def _add_high_risk_detail(self, high_risk: list) -> None:
        """高风险客户详情页（带配色表格）。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(
            slide, f"高风险客户详情",
            f"共 {len(high_risk)} 个客户需紧急干预",
        )

        if not high_risk:
            self._body_box(slide, 0.8, 1.5, 11, 1, "🎉 暂无高风险客户。", font_size=16, color=_C_GREEN)
            return

        rows = min(len(high_risk) + 1, 11)
        cols = 7
        tbl = slide.shapes.add_table(
            rows, cols, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.5 * rows),
        ).table

        # 列宽
        col_widths = [2.2, 1.5, 1.5, 1.2, 1.0, 1.5, 3.6]
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)

        headers = ["客户名称", "行业", "月 GMV", "登录天数", "投诉", "风险等级", "建议动作"]
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _C_PRIMARY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = _C_WHITE
                p.font.bold = True
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.CENTER

        actions = [
            "紧急跟进，启动挽回流程",
            "CTO 拜访 + 技术对接",
            "商务谈判，给予优惠",
            "高层拜访 + 专项支持",
            "加强 CSM 服务频次",
            "客户成功团队介入",
            "产品使用培训辅导",
            "定制解决方案",
            "客户满意度调研",
            "VIP 专属客户经理",
        ]

        for i, c in enumerate(high_risk[:10]):
            vals = [
                str(c.get("customer_name", "")),
                str(c.get("industry", "")),
                f"¥{c.get('monthly_gmv', 0):,}",
                f"{c.get('login_days', '')} 天",
                f"{c.get('complaint_count', '')} 次",
                "高",
                actions[i % len(actions)],
            ]
            for j, v in enumerate(vals):
                cell = tbl.cell(i + 1, j)
                cell.text = v
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                # 交替行颜色
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _C_WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = _C_DARK_TEXT
                    p.alignment = PP_ALIGN.CENTER

            # 风险等级列高亮
            risk_cell = tbl.cell(i + 1, 5)
            risk_cell.fill.solid()
            risk_cell.fill.fore_color.rgb = _C_RED
            for p in risk_cell.text_frame.paragraphs:
                p.font.color.rgb = _C_WHITE
                p.font.bold = True

        self._footer_line(slide)
        self._page_number(slide)

    # ---- 6. 高价值客户 ----

    def _add_high_value_slide(self, high_value: list) -> None:
        """高价值客户页（卡片式布局）。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(
            slide, "高价值客户", f"{len(high_value)} 个核心客户 — 建议配备专属 CSM",
        )

        if not high_value:
            self._body_box(slide, 0.8, 1.5, 11, 1, "暂无高价值客户。", font_size=16)
            return

        # 卡片网格：每行 4 个
        for i, c in enumerate(high_value[:8]):
            col = i % 4
            row = i // 4
            x = 0.8 + col * 3.1
            y = 1.5 + row * 2.6

            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(2.8), Inches(2.2),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _C_WHITE
            card.line.color.rgb = _C_LIGHT_GRAY
            card.line.width = Pt(0.5)

            # 左侧 accent 条
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(0.06), Inches(2.2),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = _C_GREEN
            accent.line.fill.background()

            # 客户名称
            tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.15), Inches(2.3), Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = str(c.get("customer_name", ""))
            p.font.size = Pt(15)
            p.font.color.rgb = _C_DARK_TEXT
            p.font.bold = True

            # 属性行
            attrs = [
                f"行业：{c.get('industry', '')}",
                f"月 GMV：¥{c.get('monthly_gmv', 0):,}",
                f"登录活跃：{c.get('login_days', 0)} 天",
                f"区域：{c.get('region', '')}",
            ]
            for j, attr in enumerate(attrs):
                tb2 = slide.shapes.add_textbox(
                    Inches(x + 0.25), Inches(y + 0.6 + j * 0.35), Inches(2.3), Inches(0.3),
                )
                tf2 = tb2.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = attr
                p2.font.size = Pt(11)
                p2.font.color.rgb = _C_GRAY_TEXT

        # 底部建议条
        note_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.5),
        )
        note_bar.fill.solid()
        note_bar.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
        note_bar.line.fill.background()

        tb_n = slide.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5))
        tf_n = tb_n.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = "💎  维护建议：配备专属 CSM  |  季度高层 QBR  |  提前 90 天续约  |  优先响应专属通道  |  行业标杆案例共建"
        p_n.font.size = Pt(11)
        p_n.font.color.rgb = _C_GREEN

        self._page_number(slide)

    # ---- 7. 图表网格 ----

    def _add_charts_grid(self, chart_paths: dict[str, str]) -> None:
        """行业与区域分析页（图表网格）。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(slide, "行业与区域分析", "Data Visualization")

        positions = [
            ("industry_bar", 0.3, 1.3, 6.2, 3.2),
            ("region_bar", 6.8, 1.3, 6.2, 3.2),
            ("gmv_bar", 1.5, 4.6, 10, 2.5),
        ]
        for key, left, top, width, height in positions:
            path = chart_paths.get(key)
            if path:
                slide.shapes.add_picture(
                    path, Inches(left), Inches(top), Inches(width), Inches(height),
                )

        self._page_number(slide)

    # ---- 8. AI 经营洞察总览（前言介绍页） ----

    def _add_ai_overview(self) -> None:
        """AI 经营洞察总览页：五维全景预览，图文结合，一眼看懂核心内容。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(slide, "AI 经营洞察总览", "五维全景分析 · DeepSeek 智能驱动")

        df = self._df
        total = len(df)
        high_risk = self._categories.get("high_risk", [])
        high_value = self._categories.get("high_value", [])
        growth = self._categories.get("growth", [])

        risk_counts = df["risk_level"].value_counts().to_dict()
        high_risk_n = risk_counts.get("高", 0)
        risk_pct = high_risk_n / max(total, 1) * 100
        med_risk_count = len([c for c in high_risk if "医疗" in c.get("industry", "")])
        top_complaint = max(high_risk, key=lambda c: c.get("complaint_count", 0), default={})
        top_name = top_complaint.get("customer_name", "")
        top_complaints = top_complaint.get("complaint_count", 0)

        modules = [
            {
                "title": "经营概况",
                "icon": "01",
                "color": _C_ACCENT,
                "value": str(total),
                "value_label": "总客户",
                "lines": [
                    f"覆盖 {df['industry'].nunique()} 行业 · {df['region'].nunique()} 区域",
                    f"月均 GMV ¥{df['monthly_gmv'].mean():,.0f}",
                ],
            },
            {
                "title": "风险分析",
                "icon": "02",
                "color": _C_RED,
                "value": f"{risk_pct:.0f}%",
                "value_label": "高风险占比",
                "lines": [
                    f"高风险 {high_risk_n} 个 · 中风险 {risk_counts.get('中', 0)} 个",
                    "医疗健康行业集中爆发，续约逾期严重",
                ],
            },
            {
                "title": "分层策略",
                "icon": "03",
                "color": _C_GREEN,
                "value": "3",
                "value_label": "策略层级",
                "lines": [
                    "高风险：紧急挽回 + 48h 响应",
                    "高价值：专属 CSM + QBR 复盘",
                    "增长潜力：功能激活 + 深度辅导",
                ],
            },
            {
                "title": "重点客户",
                "icon": "04",
                "color": _C_ORANGE,
                "value": str(len(high_risk)),
                "value_label": "需紧急干预",
                "lines": [
                    f"TOP 重点含 {med_risk_count} 家医疗健康",
                    f"{top_name or '—'} 投诉 {top_complaints} 次/月，优先级 P0",
                ],
            },
            {
                "title": "行动计划",
                "icon": "05",
                "color": _C_PRIMARY,
                "value": "5",
                "value_label": "本周任务",
                "lines": [
                    "P0：本周紧急回访 + 专项小组成立",
                    "P1：高价值健康度检查 + 激活计划",
                ],
            },
        ]

        # 上行 3 个卡片
        for i in range(3):
            x = 0.8 + i * 4.1
            self._render_overview_card(slide, x, 1.4, 3.5, 2.4, modules[i])

        # 下行 2 个卡片（居中）
        for i in range(2):
            x = 2.8 + i * 4.5
            self._render_overview_card(slide, x, 4.0, 3.5, 2.4, modules[3 + i])

        # 底部总结条
        summary_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.55),
        )
        summary_bar.fill.solid()
        summary_bar.fill.fore_color.rgb = _C_PRIMARY
        summary_bar.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.45))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = (
            f"▎ 核心结论：{total} 个客户中，{high_risk_n} 个处于高风险（占比 {risk_pct:.1f}%），"
            f"{len(high_value)} 个高价值客户需重点维护，"
            f"{len(growth)} 个增长潜力客户具备 uplift 空间。本周建议启动紧急挽回流程。"
        )
        p.font.size = Pt(12)
        p.font.color.rgb = _C_WHITE
        p.font.bold = True

        self._footer_line(slide)
        self._page_number(slide)

    def _render_overview_card(
        self, slide: Any, left: float, top: float,
        width: float, height: float, module: dict,
    ) -> None:
        """渲染概览卡片：圆角矩形 + 顶部色条 + 编号 + 标题 + 大数字 + 描述。"""
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _C_WHITE
        card.line.color.rgb = _C_LIGHT_GRAY
        card.line.width = Pt(0.5)

        # 顶部彩色条
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.06),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = module["color"]
        accent.line.fill.background()

        # 编号圆形
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + 0.2), Inches(top + 0.2), Inches(0.45), Inches(0.45),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = module["color"]
        circle.line.fill.background()

        tf_c = circle.text_frame
        tf_c.paragraphs[0].text = module["icon"]
        tf_c.paragraphs[0].font.size = Pt(11)
        tf_c.paragraphs[0].font.color.rgb = _C_WHITE
        tf_c.paragraphs[0].font.bold = True
        tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_c.word_wrap = False

        # 标题
        tb_t = slide.shapes.add_textbox(
            Inches(left + 0.75), Inches(top + 0.22), Inches(width - 1.0), Inches(0.4),
        )
        p_t = tb_t.text_frame.paragraphs[0]
        p_t.text = module["title"]
        p_t.font.size = Pt(16)
        p_t.font.color.rgb = _C_DARK_TEXT
        p_t.font.bold = True

        # 大数字
        tb_v = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 0.75), Inches(width - 0.4), Inches(0.55),
        )
        p_v = tb_v.text_frame.paragraphs[0]
        p_v.text = module["value"]
        p_v.font.size = Pt(28)
        p_v.font.color.rgb = module["color"]
        p_v.font.bold = True

        # 数字标签
        tb_vl = slide.shapes.add_textbox(
            Inches(left + 0.2), Inches(top + 1.25), Inches(width - 0.4), Inches(0.25),
        )
        p_vl = tb_vl.text_frame.paragraphs[0]
        p_vl.text = module["value_label"]
        p_vl.font.size = Pt(10)
        p_vl.font.color.rgb = _C_GRAY_TEXT

        # 描述行
        for j, line in enumerate(module["lines"]):
            tb_d = slide.shapes.add_textbox(
                Inches(left + 0.2), Inches(top + 1.55 + j * 0.28), Inches(width - 0.4), Inches(0.28),
            )
            p_d = tb_d.text_frame.paragraphs[0]
            p_d.text = line
            p_d.font.size = Pt(10)
            p_d.font.color.rgb = _C_DARK_TEXT

    # ---- 9. AI 经营洞察（多页，分类型视觉化渲染） ----

    def _add_ai_insights(self, insights: str) -> int:
        """
        AI 洞察多页渲染入口：解析章节 -> 合并包裹段 -> 分类型渲染。
        """
        if not insights or "（AI 洞察" in insights[:50] or "（由 AI Agent 自动生成）" in insights[:50]:
            slide = self._new_slide()
            self._page_bg(slide, _C_OFF_WHITE)
            self._title_bar(slide, "AI 经营洞察", "DeepSeek AI-Powered Analysis")
            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(1.2), Inches(0.35),
            )
            tag.fill.solid()
            tag.fill.fore_color.rgb = _C_ACCENT
            tag.line.fill.background()
            tf_tag = tag.text_frame
            p_tag = tf_tag.paragraphs[0]
            p_tag.text = "AI 分析"
            p_tag.font.size = Pt(10)
            p_tag.font.color.rgb = _C_WHITE
            p_tag.font.bold = True
            p_tag.alignment = PP_ALIGN.CENTER
            self._body_box(slide, 0.8, 2.0, 11.5, 1.5,
                           "AI 洞察内容将在配置 API Key 后自动生成。",
                           font_size=14, color=_C_GRAY_TEXT)
            self._page_number(slide)
            return 1

        sections = self._parse_insight_sections(insights)
        sections = self._merge_wrapper_sections(sections)
        if not sections:
            sections = [("AI 经营洞察", insights)]

        page_count = 0
        for idx, (title, content) in enumerate(sections):
            stype = self._detect_section_type(title)
            slide = self._new_slide()
            self._page_bg(slide, _C_OFF_WHITE)
            subtitle = f"AI 分析 \u00b7 第 {idx + 1}/{len(sections)} 部分"
            self._title_bar(slide, title, subtitle)

            if stype == "overview":
                self._render_insight_overview(slide, content)
            elif stype == "risk":
                self._render_insight_risk(slide, content)
            elif stype == "strategy":
                self._render_insight_strategy(slide, content)
            elif stype == "list":
                self._render_insight_list(slide, content)
            elif stype == "action":
                self._render_insight_action(slide, content)
            else:
                self._render_insight_default(slide, title, content)

            self._footer_line(slide)
            self._page_number(slide)
            page_count += 1

        return page_count

    # ------------------------------------------------------------------
    # AI 洞察渲染器
    # ------------------------------------------------------------------

    def _render_insight_overview(self, slide: Any, content: str) -> None:
        """经营概况：4 个 KPI 卡片 + 行业分布图 + AI 评述。"""
        import pandas as pd
        df = self._df
        kpis = [
            ("总客户数", str(len(df)), _C_ACCENT),
            ("覆盖行业", f'{df["industry"].nunique()} 个', _C_GREEN),
            ("覆盖区域", f'{df["region"].nunique()} 个', _C_ORANGE),
            ("月均 GMV", f'¥{df["monthly_gmv"].mean():,.0f}', _C_PRIMARY),
        ]
        for i, (label, value, color) in enumerate(kpis):
            self._card(slide, 0.8 + i * 3.1, 1.4, 2.8, 1.2, label, value, color)

        clean = self._clean_markdown_text(content)
        chart_path = self._chart_paths.get("industry_bar", "")
        if chart_path:
            slide.shapes.add_picture(chart_path, Inches(0.8), Inches(2.9), Inches(6.5), Inches(3.8))
            self._body_box(slide, 7.5, 2.9, 5.0, 3.8, clean[:500], font_size=11)
        else:
            self._body_box(slide, 0.8, 2.9, 11.5, 3.8, clean[:600], font_size=12)

    def _render_insight_risk(self, slide: Any, content: str) -> None:
        """风险分析：大数字 + 风险条 + 饼图 + 分析文本。"""
        df = self._df
        risk_counts = df["risk_level"].value_counts().to_dict()
        high_n = risk_counts.get("高", 0)
        mid_n = risk_counts.get("中", 0)
        low_n = risk_counts.get("低", 0)

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(4), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = str(high_n)
        p.font.size = Pt(72)
        p.font.color.rgb = _C_RED
        p.font.bold = True

        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = "高风险客户"
        p.font.size = Pt(14)
        p.font.color.rgb = _C_GRAY_TEXT

        max_n = max(high_n, mid_n, low_n, 1)
        for i, (lbl, cnt, clr) in enumerate([
            ("高", high_n, _C_RED),
            ("中", mid_n, _C_ORANGE),
            ("低", low_n, _C_GREEN),
        ]):
            y = 3.2 + i * 0.9
            tb = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(2), Inches(0.3))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{lbl}风险  {cnt}"
            p.font.size = Pt(13)
            p.font.color.rgb = clr
            p.font.bold = True
            bar_w = (cnt / max_n) * 3.5
            if bar_w > 0:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1.2), Inches(y + 0.32), Inches(bar_w), Inches(0.22),
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = clr
                bar.line.fill.background()

        chart_path = self._chart_paths.get("risk_pie", "")
        if chart_path:
            slide.shapes.add_picture(chart_path, Inches(6.5), Inches(1.4), Inches(5.5), Inches(4.5))

        clean = self._clean_markdown_text(content)
        self._body_box(slide, 0.8, 5.8, 11.5, 1.0, clean[:300], font_size=11, color=_C_GRAY_TEXT)

    def _render_insight_strategy(self, slide: Any, content: str) -> None:
        """分层经营策略：2x2 彩色策略卡片。"""
        strategies = self._extract_strategies(content)
        cards = [
            ("高风险客户挽回", _C_RED, strategies["high_risk"]),
            ("高价值客户维护", _C_GREEN, strategies["high_value"]),
            ("增长潜力客户", _C_ORANGE, strategies["growth"]),
            ("行业专项策略", _C_ACCENT, strategies["industry"]),
        ]
        positions = [(0.8, 1.5, 5.7, 2.6), (6.8, 1.5, 5.7, 2.6),
                     (0.8, 4.3, 5.7, 2.6), (6.8, 4.3, 5.7, 2.6)]

        for (label, color, desc), (x, y, w, h) in zip(cards, positions):
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _C_WHITE
            card.line.color.rgb = _C_LIGHT_GRAY
            card.line.width = Pt(0.5)

            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 0.35), Inches(y + 0.3), Inches(0.5), Inches(0.5),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(x + 1.0), Inches(y + 0.35), Inches(w - 1.3), Inches(0.4))
            p = tb.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(15)
            p.font.color.rgb = _C_DARK_TEXT
            p.font.bold = True

            tb = slide.shapes.add_textbox(Inches(x + 0.35), Inches(y + 0.95), Inches(w - 0.7), Inches(h - 1.2))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = _C_GRAY_TEXT

    def _render_insight_list(self, slide: Any, content: str) -> None:
        """重点客户清单：客户卡片（非表格），支持长名称自动换行。"""
        customers = self._categories.get("high_risk", [])
        if not customers:
            customers = self._parse_customers_from_content(content)
        if not customers:
            self._body_box(slide, 0.8, 1.5, 11, 1, "暂无重点客户数据。",
                           font_size=14, color=_C_GRAY_TEXT)
            return

        for i, c in enumerate(customers[:8]):
            col = i % 4
            row = i // 4
            x = 0.8 + col * 3.1
            y = 1.4 + row * 2.8

            name = str(c.get("customer_name", ""))
            industry = str(c.get("industry", ""))
            gmv = c.get("monthly_gmv", 0)
            login = c.get("login_days", 0)
            complaints = c.get("complaint_count", 0)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(2.4),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _C_WHITE
            card.line.color.rgb = _C_LIGHT_GRAY
            card.line.width = Pt(0.5)

            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(0.06),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = _C_RED
            accent.line.fill.background()

            tb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(2.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(13)
            p.font.color.rgb = _C_DARK_TEXT
            p.font.bold = True

            tag = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.2), Inches(y + 0.75), Inches(1.0), Inches(0.28),
            )
            tag.fill.solid()
            tag.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
            tag.line.fill.background()
            tf_tag = tag.text_frame
            tf_tag.paragraphs[0].text = industry
            tf_tag.paragraphs[0].font.size = Pt(9)
            tf_tag.paragraphs[0].font.color.rgb = _C_ACCENT
            tf_tag.paragraphs[0].alignment = PP_ALIGN.CENTER

            for j, m in enumerate([
                f"GMV: ¥{gmv:,}",
                f"登录: {login} 天/月",
                f"投诉: {complaints} 次",
            ]):
                tb = slide.shapes.add_textbox(
                    Inches(x + 0.2), Inches(y + 1.15 + j * 0.3), Inches(2.4), Inches(0.28),
                )
                p = tb.text_frame.paragraphs[0]
                p.text = m
                p.font.size = Pt(10)
                p.font.color.rgb = _C_GRAY_TEXT

    def _render_insight_action(self, slide: Any, content: str) -> None:
        """本周行动计划：横向时间轴节点。"""
        items = self._extract_action_items(content)
        if len(items) < 2:
            items = ["紧急沟通高风险客户", "启动续约谈判", "投诉处理闭环",
                     "产品培训安排", "客户健康度评估"]

        line_y = 2.5
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(line_y), Inches(11.7), Inches(0.03),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = _C_LIGHT_GRAY
        line.line.fill.background()

        colors = [_C_RED, _C_ORANGE, _C_ACCENT, _C_GREEN, _C_PRIMARY]
        max_n = min(len(items), 5)

        for i in range(max_n):
            x = 0.8 + i * 2.8
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(x + 0.6), Inches(line_y - 0.18), Inches(0.38), Inches(0.38),
            )
            node.fill.solid()
            node.fill.fore_color.rgb = colors[i]
            node.line.fill.background()

            tf = node.text_frame
            tf.paragraphs[0].text = str(i + 1)
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = _C_WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.word_wrap = False

            tb = slide.shapes.add_textbox(Inches(x), Inches(line_y + 0.4), Inches(2.4), Inches(3.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = items[i]
            p.font.size = Pt(12)
            p.font.color.rgb = _C_DARK_TEXT
            p.font.bold = True

            for sub in self._get_sub_items(i):
                p2 = tf.add_paragraph()
                p2.text = f"→ {sub}"
                p2.font.size = Pt(9)
                p2.font.color.rgb = _C_GRAY_TEXT

    def _render_insight_default(self, slide: Any, title: str, content: str) -> None:
        """默认回退渲染：纯文本布局。"""
        clean = self._clean_markdown_text(content)
        self._body_box(slide, 0.8, 1.4, 11.5, 5.5, clean[:800], font_size=12)

    def _get_sub_items(self, idx: int) -> list[str]:
        """为行动项提供默认子步骤。"""
        pool = [
            ["识别 TOP 3 风险客户", "48 小时内高管介入"],
            ["准备续约方案", "协商折扣条款"],
            ["分配投诉工单", "72 小时闭环"],
            ["确定培训时间", "准备培训材料"],
            ["收集客户反馈", "更新健康度评分"],
        ]
        return pool[idx] if idx < len(pool) else []

    # ---- 9. 行动计划 ----

    def _add_action_plan(self, high_risk: list, ai_insights: str = "") -> None:
        """行动计划页（三栏时间轴）。

        当 ai_insights 包含 AI 生成的有效内容时，尝试提取其中的行动项；
        否则使用默认模板。
        """
        slide = self._new_slide()
        self._page_bg(slide, _C_OFF_WHITE)
        self._title_bar(slide, "行动计划", "Action Plan — 分级跟进策略")

        # 尝试从 AI 洞察中提取本周行动项
        ai_actions: list[str] | None = None
        if ai_insights and not ai_insights.startswith("（AI 洞察") and "行动计划" in ai_insights:
            lines = ai_insights.split("\n")
            in_action_section = False
            collected = []
            for line in lines:
                sl = line.strip()
                if "行动计划" in sl and sl.startswith("##"):
                    in_action_section = True
                    continue
                if in_action_section:
                    if sl.startswith("##"):
                        break
                    if sl.startswith("- ") or sl.startswith("* "):
                        item = sl.lstrip("- *").strip()
                        if item:
                            collected.append(item)
            if len(collected) >= 2:
                ai_actions = collected[:6]

        # 时间轴横线
        timeline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.02),
        )
        timeline.fill.solid()
        timeline.fill.fore_color.rgb = _C_LIGHT_GRAY
        timeline.line.fill.background()

        columns = [
            ("紧急", "本周", _C_RED, 0.8, [
                f"对 {len(high_risk)} 个高风险客户启动 1 对 1 挽回",
                "CSM 主管 + 销售联合跟进逾期客户",
                "投诉客户 48 小时内完成闭环",
                "高风险客户日报追踪机制",
            ]),
            ("短期", "本月", _C_ORANGE, 5.0, [
                "高价值客户季度复盘会议",
                "增长潜力客户重点辅导计划",
                "高风险客户专项改善方案",
                "行业标杆案例内部培训",
            ]),
            ("长期", "本季度", _C_GREEN, 9.2, [
                "产品 Onboarding 流程优化",
                "客户健康度监控体系上线",
                "行业标杆案例传播推广",
                "CSM 标准化 SOP 迭代",
            ]),
        ]

        # 当从 AI 洞察中提取到有效行动项时，替换"紧急"列内容
        if ai_actions and len(ai_actions) >= 2:
            columns[0] = ("紧急（AI 建议）", "本周", _C_RED, 0.8, ai_actions[:4])

        for label, period, color, left, items in columns:
            # 时间轴节点
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(left + 1.5), Inches(1.62), Inches(0.18), Inches(0.18),
            )
            node.fill.solid()
            node.fill.fore_color.rgb = color
            node.line.fill.background()

            # 标题
            tb = slide.shapes.add_textbox(Inches(left), Inches(1.95), Inches(3.5), Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{label}  ({period})"
            p.font.size = Pt(16)
            p.font.color.rgb = color
            p.font.bold = True

            # 列表内容
            text = "\n".join(f"  {i + 1}. {item}" for i, item in enumerate(items))
            self._body_box(slide, left, 2.5, 3.8, 4.0, text, font_size=12)

        # 底部提示
        note_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.45),
        )
        note_bar.fill.solid()
        note_bar.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xE2)
        note_bar.line.fill.background()

        ai_label = "  |  AI 根据当前数据动态生成行动建议" if ai_insights and not ai_insights.startswith("（AI 洞察") else ""
        tb_n = slide.shapes.add_textbox(Inches(1.0), Inches(6.63), Inches(11.3), Inches(0.4))
        tf_n = tb_n.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = f"📋  所有任务将自动同步至 CSM 工作台  |  完成进度将在下期报告中追踪{ai_label}"
        p_n.font.size = Pt(11)
        p_n.font.color.rgb = _C_ORANGE

        self._page_number(slide)

    # ---- 10. 结束页 ----

    def _add_closing(self) -> None:
        """结束页（深色背景 + 简洁排版）。"""
        slide = self._new_slide()
        self._page_bg(slide, _C_PRIMARY)

        # 装饰线
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.0), Inches(2.3), Inches(0.03),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = _C_ACCENT
        div.line.fill.background()

        # 主文字
        tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "谢谢"
        p.font.size = Pt(52)
        p.font.color.rgb = _C_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 副文字
        tb2 = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(1.0))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "Thank You"
        p2.font.size = Pt(22)
        p2.font.color.rgb = _C_ACCENT_LIGHT
        p2.alignment = PP_ALIGN.CENTER

        # 分割线
        div2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.8), Inches(2.3), Inches(0.02),
        )
        div2.fill.solid()
        div2.fill.fore_color.rgb = _C_ACCENT
        div2.line.fill.background()

        # 来源
        tb3 = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.3), Inches(0.8))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "本报告由 AI-Customer-Success-Copilot Enterprise Reporting 自动生成"
        p3.font.size = Pt(13)
        p3.font.color.rgb = _C_GRAY_TEXT
        p3.alignment = PP_ALIGN.CENTER

        # Footer
        tb4 = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11.3), Inches(0.4))
        tf4 = tb4.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = "Multi-Tool Agent Orchestrated  |  DeepSeek AI  |  RAG Knowledge Base"
        p4.font.size = Pt(10)
        p4.font.color.rgb = _C_SECONDARY
        p4.alignment = PP_ALIGN.CENTER
